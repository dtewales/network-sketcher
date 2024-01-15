'''
SPDX-License-Identifier: Apache-2.0

Copyright 2023 Cisco Systems, Inc. and its affiliates

Licensed under the Apache License, Version 2.0 (the "License");
you may not use this file except in compliance with the License.
You may obtain a copy of the License at

http://www.apache.org/licenses/LICENSE-2.0

Unless required by applicable law or agreed to in writing, software
distributed under the License is distributed on an "AS IS" BASIS,
WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
See the License for the specific language governing permissions and
limitations under the License.
'''

from pptx import *
import sys, os, re
import numpy as np
import math
import ns_def , ns_ddx_figure
import openpyxl
from pptx import Presentation
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.dml.line import LineFormat
from pptx.shapes.connector import Connector
from pptx.util import Inches, Cm, Pt

class  ns_l3_diagram_create():
    def __init__(self):
        print('--- ns_l3_diagram_create ---')
        '''
        STEP0 get values of Master Data
        '''
        #parameter
        ws_name = 'Master_Data'
        ws_l2_name = 'Master_Data_L2'
        ws_l3_name = 'Master_Data_L3'
        excel_maseter_file = self.inFileTxt_L3_3_1.get()

        self.result_get_l2_broadcast_domains =  ns_def.get_l2_broadcast_domains.run(self,excel_maseter_file)  ## 'self.update_l2_table_array, device_l2_boradcast_domain_array, device_l2_directly_l3vport_array, device_l2_other_array, marged_l2_broadcast_group_array'

        print('--- self.update_l2_table_array ---')
        #print(self.result_get_l2_broadcast_domains[0])
        print('--- self.device_l2_boradcast_domain_array ---')
        #print(self.result_get_l2_broadcast_domains[1])
        self.device_l2_boradcast_domain_array = self.result_get_l2_broadcast_domains[1]
        print('--- device_l2_directly_l3vport_array ---')
        #print(self.result_get_l2_broadcast_domains[2])
        self.device_l2_directly_l3vport_array = self.result_get_l2_broadcast_domains[2]
        print('--- device_l2_other_array ---')
        #print(self.result_get_l2_broadcast_domains[3])
        self.device_l2_other_array = self.result_get_l2_broadcast_domains[3]
        #print('--- marged_l2_broadcast_group_array ---')
        #print(self.result_get_l2_broadcast_domains[4])
        self.marged_l2_broadcast_group_array = self.result_get_l2_broadcast_domains[4]
        print('--- self.target_l2_broadcast_group_array ---')
        print(self.target_l2_broadcast_group_array)

        self.l3_table_array = ns_def.convert_master_to_array(ws_l3_name, excel_maseter_file, '<<L3_TABLE>>')
        #print('--- self.l3_table_array  ---')
        #print(self.l3_table_array )
        self.update_l3_table_array =[]
        for tmp_l3_table_array  in self.l3_table_array :
            if tmp_l3_table_array[0] != 1 and tmp_l3_table_array[0] != 2:
                tmp_l3_table_array[1].extend(['','','',''])
                del tmp_l3_table_array[1][6:]

                #get per ip address
                work_new_l3_table_array = tmp_l3_table_array[1][4].split(',')
                add_ip_address_set_array = []
                for tmp_tmp_new_l3_table_array in work_new_l3_table_array:
                    # check IP Addresses
                    if ns_def.check_ip_format(tmp_tmp_new_l3_table_array) == 'IPv4':
                        change_tmp_ip_address = str(tmp_tmp_new_l3_table_array).replace('[', '').replace(']', '').replace('\'', '').replace(' ', '')

                        #extend [[ip_address,network_address,mask,host_address][etc..]]
                        ip_address_set_array = ns_def.get_ip_address_set(change_tmp_ip_address)
                        add_ip_address_set_array.append(ip_address_set_array)
                tmp_l3_table_array[1].append(str(add_ip_address_set_array))

                self.update_l3_table_array.append(tmp_l3_table_array[1])
        print('--- self.update_l3_table_array  ---')
        #print(self.update_l3_table_array )

        # GET way point with folder tuple
        self.wp_with_folder_tuple = {}
        for tmp_wp_folder_name in self.folder_wp_name_array[1]:
            current_row = 1
            flag_start_row = False
            flag_end_row = False

            while flag_end_row == False:
                if str(self.position_shape_tuple[current_row, 1]) == tmp_wp_folder_name:
                    start_row = current_row
                    flag_start_row = True
                if flag_start_row == True and str(self.position_shape_tuple[current_row, 1]) == '<END>':
                    flag_end_row = True
                    end_row = current_row - 1
                current_row += 1
            # print(tmp_wp_folder_name,start_row,end_row)

            for i in range(start_row, end_row + 1):
                flag_start_column = False
                current_column = 2
                while flag_start_column == False:
                    if str(self.position_shape_tuple[i, current_column]) != '<END>':
                        self.wp_with_folder_tuple[self.position_shape_tuple[i, current_column]] = tmp_wp_folder_name
                    else:
                        flag_start_column = True
                    current_column += 1

        print('---- wp_with_folder_tuple ----')
        #print(self.wp_with_folder_tuple)


        '''
        Create per area l3 ppt
        '''
        if self.click_value == 'L3-3-2':


            '''GET SIZE'''
            self.page_size_array = []
            self.slide_width = 0.0
            self.slide_hight = 0.0
            for tmp_folder_wp_name_array in self.folder_wp_name_array[0]:
                action_type = 'GET_SIZE'
                offset_x = 0.0 #inches
                offset_y = 0.0 #inches
                self.page_size_array.append( ns_l3_diagram_create.l3_area_create(self, tmp_folder_wp_name_array , action_type ,offset_x ,offset_y))

            for tmp_page_size_array in self.page_size_array:
                if self.slide_width < tmp_page_size_array[3]:
                    self.slide_width = tmp_page_size_array[3]
                if self.slide_hight < tmp_page_size_array[4]:
                    self.slide_hight = tmp_page_size_array[4]

            #add page margin
            self.slide_width += 1.0 * 2  #page margin
            self.slide_hight += 1.0 * 2  #page margin

            print('--- self.page_size_array ,self.slide_width ,self.slide_hight ---  [outline_shape_type, outline_shape_left, outline_shape_top, outline_shape_width, outline_shape_hight, folder_shape_text] ,self.slide_width ,self.slide_hight')
            #print(self.page_size_array,self.slide_width ,self.slide_hight)

            '''CREATE L3 DIAGRAM'''
            self.result_get_l2_broadcast_domains = ns_def.get_l2_broadcast_domains.run(self, excel_maseter_file)  ## 'self.update_l2_table_array, device_l2_boradcast_domain_array, device_l2_directly_l3vport_array, device_l2_other_array, marged_l2_broadcast_group_array'
            self.active_ppt = Presentation()  # define target ppt object

            for tmp_folder_wp_name_array in self.folder_wp_name_array[0]:

                action_type = 'CREATE'
                offset_x = 0.0 #inches
                offset_y = 0.0 #inches

                for tmp_page_size_array in self.page_size_array:
                    if tmp_page_size_array[5] == tmp_folder_wp_name_array:
                        offset_x = tmp_page_size_array[1]
                        offset_y = tmp_page_size_array[2]
                        break

                ns_l3_diagram_create.l3_area_create(self, tmp_folder_wp_name_array , action_type,offset_x ,offset_y)

            ### save pptx file
            self.active_ppt.save(self.output_ppt_file)


    def l3_area_create(self, target_folder_name, action_type,offset_x ,offset_y):
        print('### l3_area_create ###  ',target_folder_name,action_type,offset_x ,offset_y)
        self.used_l3segment_array = []
        ### get l3segment in the target folder
        target_all_device_array = []

        for tmp_update_l2_table_array in self.update_l2_table_array:
            if tmp_update_l2_table_array[0] == target_folder_name and 'L3' in tmp_update_l2_table_array[2]:
                #print([tmp_update_l2_table_array[1],tmp_update_l2_table_array[3]])
                target_all_device_array.append([tmp_update_l2_table_array[1],tmp_update_l2_table_array[3]])

            if tmp_update_l2_table_array[0] == target_folder_name and 'L3' in tmp_update_l2_table_array[4]:
                #print([tmp_update_l2_table_array[1], tmp_update_l2_table_array[5]])
                target_all_device_array.append([tmp_update_l2_table_array[1], tmp_update_l2_table_array[5]])

        target_all_device_array = ns_def.get_l2_broadcast_domains.get_unique_list(target_all_device_array)
        #print(target_all_device_array)

        update_l2_broadcast_group_array = []
        for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
            for tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                if tmp_tmp_target_l2_broadcast_group_array in target_all_device_array:
                    #print(tmp_tmp_target_l2_broadcast_group_array)
                    update_l2_broadcast_group_array.append(tmp_target_l2_broadcast_group_array)
                    break

        print('--- update_l2_broadcast_group_array ---')
        #print(update_l2_broadcast_group_array)

        ### get <<POSITION_SHAPE>> in MASTER EXCEL
        target_position_shape_array = []
        print('--- self.position_shape_array ---')
        #print(self.position_shape_array )

        flag_match_folder = False
        for tmp_position_shape_array in self.position_shape_array:
            if tmp_position_shape_array[1][0] == target_folder_name:
                del tmp_position_shape_array[1][0]
                del tmp_position_shape_array[1][-1]
                target_position_shape_array.append(tmp_position_shape_array[1])
                flag_match_folder = True

            elif flag_match_folder == True and tmp_position_shape_array[1][0] == '':
                del tmp_position_shape_array[1][0]
                del tmp_position_shape_array[1][-1]
                target_position_shape_array.append(tmp_position_shape_array[1])
            elif flag_match_folder == True and tmp_position_shape_array[1][0] != '':
                break

        print('--- target_position_shape_array ---')
        #print(target_position_shape_array,len(target_position_shape_array))

        ### Kyuusai if len(target_position_shape_array) == 1
        if len(target_position_shape_array) == 1:
            target_position_shape_array.append(['_AIR_9999']) # add dummy shape to second row


        ### get WP of <<POSITION_FOLDER>>in MASTER EXCEL
        wp_exist_array = [[], [], [], []]  # up/down/left/right
        target_folder_row = 1

        # check left/right wp
        for tmp_position_folder_array in self.position_folder_array:
            if tmp_position_folder_array[1][0] != '<<POSITION_FOLDER>>' and tmp_position_folder_array[1][0] != '<SET_WIDTH>':
                #print(tmp_position_folder_array)
                for index, tmp_folder_name in enumerate(tmp_position_folder_array[1]):
                    if tmp_folder_name == target_folder_name:
                        target_folder_row = tmp_position_folder_array[0]
                        if len(tmp_position_folder_array[1]) >= (index + 2):
                            if '_wp_' in str(tmp_position_folder_array[1][index+1]):
                                #print('right ', tmp_position_folder_array[1][index+1])
                                wp_exist_array[3].extend([tmp_position_folder_array[1][index+1]])
                        if index >= 1:
                            if '_wp_' in str(tmp_position_folder_array[1][index - 1]):
                                #print('left ', tmp_position_folder_array[1][index - 1])
                                wp_exist_array[2].extend([tmp_position_folder_array[1][index - 1]])

        #check up/down wp
        for tmp_position_folder_array in self.position_folder_array:
            if tmp_position_folder_array[0]  == target_folder_row -2:
                for tmp_tmp_position_folder_array in tmp_position_folder_array[1]:
                    if '_wp_' in str(tmp_tmp_position_folder_array):
                        wp_exist_array[0].extend([tmp_tmp_position_folder_array])

            if tmp_position_folder_array[0]  == target_folder_row +2:
                for tmp_tmp_position_folder_array in tmp_position_folder_array[1]:
                    if '_wp_' in str(tmp_tmp_position_folder_array):
                        wp_exist_array[1].extend([tmp_tmp_position_folder_array])

        #print('--- wp_exist_array ---')
        #print(wp_exist_array,'  #up/down/left/right')

        #convert _wp_ folder name to shape name
        new_wp_exist_array = [[], [], [], []]

        for index,tmp_wp_exist_array in enumerate(wp_exist_array):
            for tmp_tmp_wp_exist_array in tmp_wp_exist_array:
                for tmp_wp_with_folder_tuple in self.wp_with_folder_tuple:
                    if tmp_tmp_wp_exist_array == str(self.wp_with_folder_tuple[tmp_wp_with_folder_tuple]):
                        #check if wp is connected to own folder

                        #check wp name
                        for tmp_update_l2_broadcast_group_array in update_l2_broadcast_group_array:
                            for tmp_tmp_update_l2_broadcast_group_array in tmp_update_l2_broadcast_group_array[1]:
                                #print(tmp_tmp_target_l2_broadcast_group_array )
                                if tmp_tmp_update_l2_broadcast_group_array[0] == str(tmp_wp_with_folder_tuple):
                                    # print(index,tmp_wp_with_folder_tuple)
                                    new_wp_exist_array[index].extend([tmp_wp_with_folder_tuple])

                        #check broadcast domain number
                        for tmp_device_l2_boradcast_domain_array in self.device_l2_boradcast_domain_array:
                            for tmp_tmp_device_l2_boradcast_domain_array in tmp_device_l2_boradcast_domain_array:
                                #print(tmp_tmp_device_l2_boradcast_domain_array)
                                if tmp_tmp_device_l2_boradcast_domain_array[1] == str(tmp_wp_with_folder_tuple):
                                    for tmp_update_l2_broadcast_group_array in update_l2_broadcast_group_array:
                                        for tmp_tmp_update_l2_broadcast_group_array in tmp_update_l2_broadcast_group_array:
                                            if  tmp_tmp_device_l2_boradcast_domain_array[0] in tmp_tmp_update_l2_broadcast_group_array :
                                                #print(tmp_tmp_device_l2_boradcast_domain_array[0],tmp_tmp_device_l2_boradcast_domain_array[1])
                                                new_wp_exist_array[index].extend([tmp_wp_with_folder_tuple])

        new_wp_exist_array = [list(set(new_wp_exist_array [0])),list(set(new_wp_exist_array [1])),list(set(new_wp_exist_array [2])),list(set(new_wp_exist_array[3]))]

        print('--- new_wp_exist_array ---')
        #print(new_wp_exist_array ,'  #up/down/left/right')
        self.new_wp_exist_array = new_wp_exist_array

        marge_target_position_shape_array = target_position_shape_array
        if new_wp_exist_array[0] != []:
            marge_target_position_shape_array.insert(0,new_wp_exist_array[0])
        if new_wp_exist_array[1] != []:
            marge_target_position_shape_array.append(new_wp_exist_array[1])

        print('--- marge_target_position_shape_array ---')
        #print(marge_target_position_shape_array)

        wp_marge_target_position_shape_array = marge_target_position_shape_array

        #left_right_wp_row_num = math.floor(len(marge_target_position_shape_array) * 0.5 - 1)
        left_right_wp_row_num = math.floor((len(marge_target_position_shape_array) - 1) * 0.5)     # updated

        # write wp left
        if new_wp_exist_array[2] != []:
            for tmp_i in new_wp_exist_array[2]:
                shape_text = tmp_i

                # insert wp to edge of left side
                wp_marge_target_position_shape_array[left_right_wp_row_num].insert(0, shape_text)

        # write wp right
        if new_wp_exist_array[3] != []:
            for tmp_i in new_wp_exist_array[3]:
                shape_text = tmp_i

                # insert wp to edge of right side
                wp_marge_target_position_shape_array[left_right_wp_row_num].extend([shape_text])

        print('--- wp_marge_target_position_shape_array---')
        #print(wp_marge_target_position_shape_array)

        #get up down l3if array
        self.up_down_l3if_array = get_up_down_l3if_count(self,wp_marge_target_position_shape_array)

        print('--- self.up_down_l3if_array ---  # up/down')
        #print(self.up_down_l3if_array[0])
        #print(self.up_down_l3if_array[1])

        ### get index for target_position_shape_array
        #get index_1  -> target_position_shape_array
        self.index_1_array = []
        for self.index_1,tmp_target_position_shape_array in enumerate(target_position_shape_array):
            #print(self.index_1,tmp_target_position_shape_array )
            self.index_1_array.append(self.index_1)

        #get index_11  -> marge_target_position_shape_array
        self.index_11_array = []
        for self.index_11,tmp_target_position_shape_array in enumerate(marge_target_position_shape_array):
            #print(self.index_11,tmp_target_position_shape_array )
            self.index_11_array.append(self.index_1)


        ### GET identify L3IF that is connected to l3 segment
        self.l3_if_has_l3_segment_array = []
        return_get_l3_segment_array = []

        self.used_l3segment_array = []
        for self.index_2, tmp_target_position_shape_array in enumerate(target_position_shape_array):
            for return_get_l3_segment_num in get_l3_segment_num(self,tmp_target_position_shape_array,target_position_shape_array)[1]:
                return_get_l3_segment_array.append(return_get_l3_segment_num)

                for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
                    if return_get_l3_segment_num in tmp_target_l2_broadcast_group_array[1]:
                        self.l3_if_has_l3_segment_array.extend(tmp_target_l2_broadcast_group_array[1])

        print('--- self.l3_if_has_l3_segment_array ---')
        #print(self.l3_if_has_l3_segment_array)

        ### GET L3 instance
        self.defalut_l3_instance_name = 'Defalut'
        self.l3_instance_array = []
        self.update_l3_instance_array = []
        for tmp_update_l3_table_array in self.update_l3_table_array:
            if tmp_update_l3_table_array[3].replace(' ','') != '':
                self.l3_instance_array.append([tmp_update_l3_table_array[1],tmp_update_l3_table_array[3]])

        self.l3_instance_array = ns_def.get_l2_broadcast_domains.get_unique_list(self.l3_instance_array)
        print('--- self.l3_instance_array ---')
        #print(self.l3_instance_array)

        for tmp_update_l3_table_array in self.update_l3_table_array:
            for tmp_l3_instance_array  in self.l3_instance_array:
                if tmp_l3_instance_array[0] == tmp_update_l3_table_array[1]:
                    if tmp_update_l3_table_array[3].replace(' ','') != '':
                        self.update_l3_instance_array.append([tmp_update_l3_table_array,tmp_update_l3_table_array[3]])
                        #print([tmp_update_l3_table_array,tmp_update_l3_table_array[3]])
                    else:
                        self.update_l3_instance_array.append([tmp_update_l3_table_array, self.defalut_l3_instance_name])
                        #print([tmp_update_l3_table_array, self.defalut_l3_instance_name])

        print('--- self.update_l3_instance_array ---')
        #print(self.update_l3_instance_array)


        '''
        make ppt diagram
        '''
        ### create slide
        if action_type == 'CREATE':
            self.active_ppt.slide_width = Inches(10.0)
            self.active_ppt.slide_height = Inches(5.0)

            ### adjust to large size , ver 2.1
            if self.slide_width > 56.0:
                self.slide_width = 56.0
            if self.slide_hight > 56.0:
                self.slide_hight = 56.0

            #input from get_size
            if self.active_ppt.slide_width < Inches(self.slide_width):
                self.active_ppt.slide_width = Inches(self.slide_width)
            if self.active_ppt.slide_height < Inches(self.slide_hight):
                self.active_ppt.slide_height = Inches(self.slide_hight)

            self.title_only_slide_layout = self.active_ppt.slide_layouts[5]
            self.slide = self.active_ppt.slides.add_slide(self.title_only_slide_layout)
            self.slide.shapes.title.left = Inches(0.0)
            self.slide.shapes.title.top = Inches(0.0)
            self.slide.shapes.title.width = Inches(10.0)
            self.slide.shapes.title.height = Inches(1.0)
            self.shape = self.slide.shapes

            self.shape.title.text = '[L3] ' + target_folder_name

        ### parameter
        self.left_margin = 1.0 # Inches
        self.top_margin = 1.0 # Inches

        # input from get_size
        if action_type == 'CREATE':
            self.left_margin = self.left_margin - offset_x + self.left_margin # Inches
            self.top_margin = self.top_margin - offset_y + self.top_margin  # Inches

        top_offset = 0.0 #Inches
        left_offset = 0.0 #Inches

        self.folder_font_type = 'Calibri'
        self.folder_font_size = 10  # Pt
        self.shape_font_type = 'Calibri'
        self.shae_font_size = 6.0  # Pt
        self.shae_font_large_size = 8.0  # Pt
        self.tag_font_large_size = 4.0 # Pt

        self.between_shape_column = 0.5 #inches
        between_shape_row = 0.25  # inches
        self.between_l3if = 0.25  #inches
        l3_segment_up_down_offset = 0.15 # inches
        min_between_line = 0.075  # inches
        min_shape_width = 1.0 #inches

        l3_segment_hight_ratio = 1.75 # ratio


        ''' 
        main loop 
        '''
        shape_left = self.left_margin
        shape_top = self.top_margin

        '''
        loop write device
        '''
        l3segment_line_array = []
        self.connected_l3if_key_array = []
        self.all_l3if_tag_array = []
        self.mark_multi_ip_array = []
        self.size_l3_instance_array = []
        self.used_l3segment_array = []
        self.area_position_array = [999.0,0.0,0.0,0.0,target_folder_name]  # shape_left, shape_top, shape_width, shape_hight, shape_text
        self.outline_position_array = [999.0, 0.0, 0.0, 0.0]  # shape_left, shape_top, shape_width, shape_hight
        max_offset_x = 0.0
        end_l3_seg_inche_x = 0.0
        self.mark_wp_top = self.top_margin + top_offset
        flag_first_colmun = True

        for self.index_2,tmp_target_position_shape_array in enumerate(target_position_shape_array):
            start_l3_seg_inche_x = self.left_margin + left_offset

            ''' write device and wp(up/down)'''
            for tmp_tmp_target_position_shape_array in tmp_target_position_shape_array:
                tmp_left_array = []
                tmp_right_array = []
                for i in new_wp_exist_array[2]:
                    tmp_left_array.append(i)
                for k in new_wp_exist_array[3]:
                    tmp_left_array.append(k)


                if tmp_tmp_target_position_shape_array not in tmp_left_array  and tmp_tmp_target_position_shape_array not in tmp_right_array and '_AIR_' not in tmp_tmp_target_position_shape_array:  # except left/right wp in writing pre device. If you need _AIR_ empty space, delete '_AIR_ not in tmp_tmp_target_position_shape_array'
                    shape_text  = tmp_tmp_target_position_shape_array
                    shape_type  = 'DEVICE_NORMAL'
                    shape_left  = self.left_margin + left_offset
                    shape_top   = self.top_margin + top_offset
                    shape_width_hight_array = ns_def.get_description_width_hight(self.shae_font_large_size,tmp_tmp_target_position_shape_array)
                    shape_width = shape_width_hight_array[0]
                    shape_hight = shape_width_hight_array[1] * 5

                    if shape_text in self.wp_list_array:
                        shape_type = 'WAY_POINT_NORMAL'

                    self.shape_width_if_array = get_shape_width_if_array(self,tmp_tmp_target_position_shape_array)  ### return_shape_width,tmp_up_array,tmp_down_array,full_ip_address_width_array
                    tmp_shpae_width = self.shape_width_if_array[0]
                    if shape_width < tmp_shpae_width:
                        shape_width = tmp_shpae_width

                    if min_shape_width > shape_width:
                        shape_width = min_shape_width

                    #check l3 instance exist
                    for tmp_l3_instance_array in self.l3_instance_array:
                        if tmp_l3_instance_array[0] == shape_text:
                            shape_type = 'DEVICE_L3_INSTANCE'

                    #write l3 instance
                    self.between_l3instance = min_between_line * 2
                    tmp_l3_add_shape_array = []
                    if shape_type == 'DEVICE_L3_INSTANCE':
                        tmp_l3_instance_array = []
                        for tmp_update_l3_instance_array in self.update_l3_instance_array:
                            if tmp_update_l3_instance_array[0][1] == shape_text:
                                tmp_l3_instance_array.append(tmp_update_l3_instance_array[1])

                        tmp_l3_instance_array = ns_def.get_l2_broadcast_domains.get_unique_list(tmp_l3_instance_array)
                        offset_l3_instance = ns_def.get_description_width_hight(self.shae_font_size,shape_text)[0] + self.between_l3instance

                        calc_need_device_width = ns_def.get_description_width_hight(self.shae_font_size,shape_text)[0] + offset_l3_instance
                        for tmp_tmp_l3_instance_array in tmp_l3_instance_array:
                            l3_shape_text = tmp_tmp_l3_instance_array
                            l3_shape_width = ns_def.get_description_width_hight(self.shae_font_size,tmp_tmp_l3_instance_array)[0]
                            l3_shape_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_tmp_l3_instance_array)[1]

                            if (min_shape_width * 0.3) > l3_shape_width:
                                l3_shape_width = (min_shape_width * 0.3)

                            l3_shape_top = shape_top + (shape_hight * 0.5) - l3_shape_hight * 0.5
                            l3_shape_left = shape_left + offset_l3_instance
                            l3_shape_type = 'L3_INSTANCE'

                            calc_need_device_width += l3_shape_width + self.between_l3instance

                            if action_type == 'CREATE':
                                #self.shape = self.slide.shapes
                                #ns_ddx_figure.extended.add_shape(self, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight, l3_shape_text)
                                tmp_l3_add_shape_array.append([l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight, l3_shape_text])

                            offset_l3_instance += l3_shape_width + self.between_l3instance
                            #self.size_l3_instance_array.append([shape_text, l3_shape_text, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight])

                        if calc_need_device_width > shape_width:
                            ### extend distance of shape

                            shape_width = calc_need_device_width
                            for tmp_tmp_l3_add_shape_array in tmp_l3_add_shape_array:
                                l3_shape_type = tmp_tmp_l3_add_shape_array[0]
                                l3_shape_left = tmp_tmp_l3_add_shape_array[1]
                                l3_shape_top = tmp_tmp_l3_add_shape_array[2]
                                l3_shape_width = tmp_tmp_l3_add_shape_array[3]
                                l3_shape_hight = tmp_tmp_l3_add_shape_array[4]
                                l3_shape_text = tmp_tmp_l3_add_shape_array[5]
                                self.shape = self.slide.shapes
                                ns_ddx_figure.extended.add_shape(self, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight, l3_shape_text)

                                self.size_l3_instance_array.append([shape_text, l3_shape_text, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight])

                        else:
                            ### not extend distance of shape
                            tmp_l3_instance_array = []
                            for tmp_update_l3_instance_array in self.update_l3_instance_array:
                                if tmp_update_l3_instance_array[0][1] == shape_text:
                                    tmp_l3_instance_array.append(tmp_update_l3_instance_array[1])

                            tmp_l3_instance_array = ns_def.get_l2_broadcast_domains.get_unique_list(tmp_l3_instance_array)
                            tmp_plus_width = shape_width / (len(tmp_l3_instance_array) + 1)
                            offset_l3_instance = 0.0

                            for tmp_tmp_l3_instance_array in tmp_l3_instance_array:
                                l3_shape_text = tmp_tmp_l3_instance_array
                                l3_shape_width = ns_def.get_description_width_hight(self.shae_font_size, tmp_tmp_l3_instance_array)[0]
                                l3_shape_hight = ns_def.get_description_width_hight(self.shae_font_size, tmp_tmp_l3_instance_array)[1]

                                if (min_shape_width * 0.3) > l3_shape_width:
                                    l3_shape_width = (min_shape_width * 0.3)

                                l3_shape_top = shape_top + (shape_hight * 0.5) - l3_shape_hight * 0.5
                                l3_shape_left = shape_left + tmp_plus_width - l3_shape_width * 0.5 + offset_l3_instance
                                l3_shape_type = 'L3_INSTANCE'

                                if action_type == 'CREATE':
                                    self.shape = self.slide.shapes
                                    ns_ddx_figure.extended.add_shape(self, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight, l3_shape_text)

                                offset_l3_instance += tmp_plus_width
                                self.size_l3_instance_array.append([shape_text, l3_shape_text, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight])


                    if '_AIR_' not in shape_text:
                        if action_type == 'CREATE':
                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)
                            self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
                            self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer

                        '''GET Folder and Outline position'''
                        # get folder left
                        if (shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL' or shape_type == 'WAY_POINT_NORMAL') and self.area_position_array[0] + self.between_shape_column  > shape_left:
                            self.area_position_array[0] = shape_left - self.between_shape_column
                            self.outline_position_array[0] = shape_left - self.between_shape_column * 2

                        # get folder top
                        if (shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL') and flag_first_colmun == True:
                            self.area_position_array[1] = shape_top - between_shape_row
                            if new_wp_exist_array[0] == []:
                                self.outline_position_array[1] = shape_top - between_shape_row * 2
                            flag_first_colmun = False

                        # get folder width
                        if (shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL' or shape_type == 'WAY_POINT_NORMAL') and (self.outline_position_array[0] + self.outline_position_array[2]) < shape_left + shape_width:
                            self.outline_position_array[2] = shape_left + shape_width + self.between_shape_column * 2 - self.outline_position_array[0]

                        if (shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL' ) and (self.area_position_array[0] + self.area_position_array[2]) < shape_left + shape_width:
                            self.area_position_array[2] = shape_left + shape_width + self.between_shape_column - self.area_position_array[0]

                        # get folder hight
                        if shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL':
                            self.area_position_array[3] = shape_top + shape_hight + between_shape_row - self.area_position_array[1]
                            self.outline_position_array[3] = shape_top + shape_hight + between_shape_row * 3 - self.outline_position_array[1]

                        # adjust shape_hight if downside wp exist
                        if new_wp_exist_array[1] != []: #up/down/left/right
                            self.area_position_array[3] = (shape_top - between_shape_row) - self.area_position_array[1]


                        #adjust shape top of outline, if wp
                        if shape_type == 'WAY_POINT_NORMAL':
                            if shape_text in  new_wp_exist_array[0]:
                                self.outline_position_array[1] = shape_top - between_shape_row

                            if shape_text in new_wp_exist_array[1]:
                                self.outline_position_array[3] = (shape_top + shape_hight + between_shape_row) - self.outline_position_array[1]


                    left_offset += shape_width + self.between_shape_column

                    #print('--- self.size_l3_instance_array ---  shape_text, l3_shape_text, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight')
                    #print(self.size_l3_instance_array)

                    '''write l3 if '''

                    #print('### self.shape_width_if_array[1], self.shape_width_if_array[2]  ',self.shape_width_if_array[1], self.shape_width_if_array[2])
                    tag_up_offset_x = self.between_l3if
                    tag_down_offset_x = self.between_l3if
                    for tmp_update_l3_table_array  in self.update_l3_table_array:
                        if tmp_update_l3_table_array[1] == shape_text:
                            '''write upside l3 if'''
                            for up_shape_width_if_array in self.shape_width_if_array[1]:
                                if up_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                    #print('##UP   ',shape_text,up_shape_width_if_array[1])
                                    shape_width_hight_array = ns_def.get_description_width_hight(self.tag_font_large_size, up_shape_width_if_array[1]) # width, hight
                                    tag_shape_type  = 'TAG_NORMAL'
                                    tag_shape_left  = shape_left + tag_up_offset_x
                                    tag_shape_top   = shape_top - shape_width_hight_array[1] * 0.5
                                    tag_shape_width = shape_width_hight_array[0]
                                    tag_shape_hight = shape_width_hight_array[1]
                                    tag_shape_text  = up_shape_width_if_array[1]

                                    if action_type == 'CREATE':
                                        self.shape = self.slide.shapes
                                        ns_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                    self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,up_shape_width_if_array[0],shape_text])

                                    #reflect description ip address name distance-1
                                    flag_match_shape_width_if_array = False
                                    for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                        if [shape_text ,up_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                            tag_up_offset_x += tmp_shape_width_if_array[1]
                                            flag_match_shape_width_if_array = True
                                            tmp_add_width = tmp_shape_width_if_array[1]

                                    if flag_match_shape_width_if_array == False:
                                        tag_up_offset_x += tag_shape_width + self.between_l3if
                                        tmp_add_width = tag_shape_width + self.between_l3if

                                    '''mark ip address(up side)'''
                                    offset_ipaddress = 0.0 #inches
                                    remake_array = []
                                    remake_array = eval(tmp_update_l3_table_array[6])
                                    if remake_array != []:
                                        #print('##mark ip address', len(remake_array),remake_array)
                                        for tmp_remake_array in remake_array:
                                            tag_shape_type = 'IP_ADDRESS_TAG'
                                            tag_ip_width = ns_def.get_description_width_hight(self.shae_font_size,tmp_remake_array[2])[0]
                                            self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2],tmp_remake_array,shape_text,len(remake_array),[shape_text,up_shape_width_if_array[0]],tag_shape_left ])

                                            offset_ipaddress += tag_shape_hight

                                            # reflect description ip address name distance-2
                                            #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                            #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width


                            #write downside l3 if
                            for down_shape_width_if_array in self.shape_width_if_array[2]:
                                if down_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                    #print('##DOWN ',shape_text,down_shape_width_if_array[1])
                                    shape_width_hight_array = ns_def.get_description_width_hight(self.tag_font_large_size, down_shape_width_if_array[1]) # width, hight
                                    tag_shape_type  = 'TAG_NORMAL'
                                    tag_shape_left  = shape_left + tag_down_offset_x
                                    tag_shape_top   = shape_top + shape_hight  - shape_width_hight_array[1] * 0.5
                                    tag_shape_width = shape_width_hight_array[0]
                                    tag_shape_hight = shape_width_hight_array[1]
                                    tag_shape_text  = down_shape_width_if_array[1]

                                    if action_type == 'CREATE':
                                        self.shape = self.slide.shapes
                                        ns_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                    self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,down_shape_width_if_array[0],shape_text])

                                    # reflect description ip address name distance-1
                                    flag_match_shape_width_if_array = False
                                    for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                        if [shape_text, down_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                            tag_down_offset_x += tmp_shape_width_if_array[1]
                                            flag_match_shape_width_if_array = True
                                            tmp_add_width = tmp_shape_width_if_array[1]

                                    if flag_match_shape_width_if_array == False:
                                        tag_down_offset_x += tag_shape_width + self.between_l3if
                                        tmp_add_width = tag_shape_width + self.between_l3if

                                    '''mark ip address(down side)'''
                                    offset_ipaddress = 0.0 #inches
                                    remake_array = []
                                    remake_array = eval(tmp_update_l3_table_array[6])
                                    if remake_array != []:
                                        #print('##mark ip address', len(remake_array),remake_array)
                                        for tmp_remake_array in remake_array:
                                            tag_shape_type = 'IP_ADDRESS_TAG'
                                            tag_ip_width = ns_def.get_description_width_hight(self.shae_font_size,tmp_remake_array[2])[0]
                                            self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top + tag_shape_hight + offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2],tmp_remake_array,shape_text,len(remake_array),[shape_text,down_shape_width_if_array[0]],tag_shape_left])

                                            offset_ipaddress += tag_shape_hight

                                            # reflect description ip address name distance-2
                                            #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                            #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width

            ''' pre mark wp(left/right)'''
            if (new_wp_exist_array[2] != [] or new_wp_exist_array[3] != []) and left_right_wp_row_num == self.index_2:
                self.mark_wp_top = self.top_margin + top_offset

            '''add top_offset for folder outline'''
            #upside
            if shape_text in self.wp_list_array:
                top_offset += between_shape_row * 2

            #downside
            if len(target_position_shape_array) > self.index_2 + 1:
                for down_target_position_shape_array in target_position_shape_array[self.index_2 + 1]:
                    if down_target_position_shape_array in self.wp_list_array and len(target_position_shape_array) == self.index_2 +2:
                        top_offset += between_shape_row * 2
                        break

            ### check end l3 segment point
            if end_l3_seg_inche_x < shape_left + shape_width:
                end_l3_seg_inche_x = shape_left + shape_width

            ### write broadcast domain line
            top_device_name_array = tmp_target_position_shape_array

            #get count l3_segment
            return_get_l3_segment_num = get_l3_segment_num(self,top_device_name_array ,target_position_shape_array)
            count_l3segment = return_get_l3_segment_num[0]
            self.connected_l3if_key_array.append(return_get_l3_segment_num[1])
            #print('##  return_get_l3_segment_num ',return_get_l3_segment_num ,top_device_name_array )

            tmp_l3segment_y_array = []
            if count_l3segment != 0:
                tmp_line_offset = l3_segment_up_down_offset
                for tmp_count_l3segment in range(count_l3segment):
                    #add distance upside or downside of device shape
                    if tmp_count_l3segment == 0:
                        top_offset += between_shape_row + l3_segment_up_down_offset
                    else:
                        top_offset += between_shape_row
                    if tmp_count_l3segment +1 == count_l3segment:
                        top_offset += l3_segment_up_down_offset

                    tmp_l3segment_y_array.append(shape_top + tmp_line_offset  + shape_hight + between_shape_row)
                    tmp_line_offset += between_shape_row
            l3segment_line_array.append([[start_l3_seg_inche_x,end_l3_seg_inche_x],tmp_l3segment_y_array,return_get_l3_segment_num[1]])

            '''change offset  check_move_to_right '''
            top_offset += shape_hight + between_shape_row

            now_offset_x = end_l3_seg_inche_x - self.left_margin
            if max_offset_x < now_offset_x :
                max_offset_x = now_offset_x

            if check_move_to_right(self,top_device_name_array,target_position_shape_array) == True:
                left_offset = max_offset_x
            else:
                left_offset = start_l3_seg_inche_x - self.left_margin

        #print('## end_l3_seg_inche_x ', end_l3_seg_inche_x)
        print('--- l3segment_line_array ---   [start_l3_seg_inche_x,end_l3_seg_inche_x],tmp_l3segment_y_array, return_get_l3_segment_num[1]')
        #print(l3segment_line_array)


        ''' 
        write wp(left/right)
        '''
        shape_top = self.mark_wp_top
        shape_width_hight_array = ns_def.get_description_width_hight(self.shae_font_large_size, str(new_wp_exist_array[2]))
        shape_width = shape_width_hight_array[0]
        shape_hight = shape_width_hight_array[1] * 5


        ''' write wp left'''
        offset_shape_left = 0.0
        if new_wp_exist_array[2] != []:
            for tmp_i in new_wp_exist_array[2]:
                shape_text = tmp_i
                shape_type = 'WAY_POINT_NORMAL'

                self.shape_width_if_array = get_shape_width_if_array(self, shape_text)  # return_shape_width, tmp_up_array, tmp_down_array
                tmp_shpae_width = self.shape_width_if_array[0]
                if shape_width < tmp_shpae_width:
                    shape_width = tmp_shpae_width

                if min_shape_width > shape_width:
                    shape_width = min_shape_width

                print('### WRITE LEFT WP  ', shape_text, shape_width,new_wp_exist_array[2])
                shape_left = self.left_margin - shape_width - self.between_shape_column * 3 + offset_shape_left

                if action_type == 'CREATE':
                    self.shape = self.slide.shapes
                    ns_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)

                #get left side folder and outline point
                self.area_position_array[0] = shape_left + shape_width + self.between_shape_column * 2 - offset_shape_left

                self.outline_position_array[0] = shape_left - self.between_shape_column * 2
                self.outline_position_array[2] += shape_width + self.between_shape_column * 3

                '''write wp_left l3 if '''
                # print('### self.shape_width_if_array[1], self.shape_width_if_array[2]  ',self.shape_width_if_array[1], self.shape_width_if_array[2])
                tag_up_offset_x = self.between_l3if
                tag_down_offset_x = self.between_l3if

                for tmp_update_l3_table_array in self.update_l3_table_array:
                    if tmp_update_l3_table_array[1] == shape_text:
                        for up_shape_width_if_array in self.shape_width_if_array[1]:
                            #write up side l3 if
                            if up_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                # print('##UP   ',shape_text,up_shape_width_if_array[1])
                                shape_width_hight_array = ns_def.get_description_width_hight(self.tag_font_large_size, up_shape_width_if_array[1])  # width, hight
                                tag_shape_type = 'TAG_NORMAL'
                                tag_shape_left = shape_left + tag_up_offset_x
                                tag_shape_top = shape_top - shape_width_hight_array[1] * 0.5
                                tag_shape_width = shape_width_hight_array[0]
                                tag_shape_hight = shape_width_hight_array[1]
                                tag_shape_text = up_shape_width_if_array[1]

                                if action_type == 'CREATE':
                                    self.shape = self.slide.shapes
                                    ns_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,up_shape_width_if_array[0],shape_text])

                                # reflect description ip address name distance-1
                                flag_match_shape_width_if_array = False
                                for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                    if [shape_text, up_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                        tag_up_offset_x += tmp_shape_width_if_array[1]
                                        flag_match_shape_width_if_array = True
                                        tmp_add_width = tmp_shape_width_if_array[1]

                                if flag_match_shape_width_if_array == False:
                                    tag_up_offset_x += tag_shape_width + self.between_l3if
                                    tmp_add_width = tag_shape_width + self.between_l3if


                                '''wp_left write ip address(up side)'''
                                offset_ipaddress = 0.0  # inches
                                remake_array = []
                                remake_array = eval(tmp_update_l3_table_array[6])
                                if remake_array != []:
                                    for tmp_remake_array in remake_array:
                                        #print('##wp_left write ip address', tmp_remake_array)
                                        tag_shape_type = 'IP_ADDRESS_TAG'
                                        tag_ip_width = ns_def.get_description_width_hight(self.shae_font_size, tmp_remake_array[2])[0]
                                        #self.shape = self.slide.shapes

                                        self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2], tmp_remake_array, shape_text, len(remake_array), [shape_text, up_shape_width_if_array[0]], tag_shape_left])
                                        offset_ipaddress += tag_shape_hight

                                        # reflect description ip address name distance-2
                                        #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                        #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width

                        # wite down side l3 if
                        for down_shape_width_if_array in self.shape_width_if_array[2]:
                            if down_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                #print('##DOWN ',shape_text,down_shape_width_if_array[1])
                                shape_width_hight_array = ns_def.get_description_width_hight(self.tag_font_large_size, down_shape_width_if_array[1]) # width, hight
                                tag_shape_type  = 'TAG_NORMAL'
                                tag_shape_left  = shape_left + tag_down_offset_x
                                tag_shape_top   = shape_top + shape_hight  - shape_width_hight_array[1] * 0.5
                                tag_shape_width = shape_width_hight_array[0]
                                tag_shape_hight = shape_width_hight_array[1]
                                tag_shape_text  = down_shape_width_if_array[1]

                                if action_type == 'CREATE':
                                    self.shape = self.slide.shapes
                                    ns_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,down_shape_width_if_array[0],shape_text])

                                # reflect description ip address name distance-1
                                flag_match_shape_width_if_array = False
                                for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                    if [shape_text, down_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                        tag_up_offset_x += tmp_shape_width_if_array[1]
                                        flag_match_shape_width_if_array = True
                                        tmp_add_width = tmp_shape_width_if_array[1]

                                if flag_match_shape_width_if_array == False:
                                    tag_down_offset_x += tag_shape_width + self.between_l3if
                                    tmp_add_width = tag_shape_width + self.between_l3if

                                '''wp_left write ip address(down side)'''
                                offset_ipaddress = 0.0  # inches
                                remake_array = []
                                remake_array = eval(tmp_update_l3_table_array[6])
                                if remake_array != []:
                                    for tmp_remake_array in remake_array:
                                        #print('##wp_left write ip address', tmp_remake_array)
                                        tag_shape_type = 'IP_ADDRESS_TAG'
                                        tag_ip_width = ns_def.get_description_width_hight(self.shae_font_size, tmp_remake_array[2])[0]

                                        #self.shape = self.slide.shapes
                                        self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top + tag_shape_hight + offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2], tmp_remake_array, shape_text, len(remake_array), [shape_text, down_shape_width_if_array[0]],tag_shape_left])

                                        offset_ipaddress += tag_shape_hight

                                        # reflect description ip address name distance-2
                                        #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                        #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width

                offset_shape_left = offset_shape_left - shape_width - self.between_shape_column

        ''' write wp right'''
        offset_shape_right = 0.0
        if new_wp_exist_array[3] != []:

            for tmp_i in reversed(new_wp_exist_array[3]):
                shape_text = tmp_i

                shape_type = 'WAY_POINT_NORMAL'
                self.shape_width_if_array = get_shape_width_if_array(self, shape_text)  # return_shape_width, tmp_up_array, tmp_down_array

                tmp_shpae_width = self.shape_width_if_array[0]

                if shape_width < tmp_shpae_width:
                    shape_width = tmp_shpae_width

                if min_shape_width > shape_width:
                    shape_width = min_shape_width

                print('### WRITE RIGHT WP  ', new_wp_exist_array[3], shape_width)
                #shape_left = end_l3_seg_inche_x + self.between_shape_column * 3
                shape_left = self.area_position_array[0] + self.area_position_array[2] + self.between_shape_column * 2  + offset_shape_right # updated
                if action_type == 'CREATE':
                    self.shape = self.slide.shapes
                    ns_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)

                #get left side folder and outline point
                self.outline_position_array[2] += shape_width + self.between_shape_column * 3

                '''write wp_right l3 if '''
                tag_up_offset_x = self.between_l3if
                tag_down_offset_x = self.between_l3if

                for tmp_update_l3_table_array in self.update_l3_table_array:
                    if tmp_update_l3_table_array[1] == shape_text:
                        for up_shape_width_if_array in self.shape_width_if_array[1]:
                            # write up side l3 if
                            if up_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                # print('##UP   ',shape_text,up_shape_width_if_array[1])
                                shape_width_hight_array = ns_def.get_description_width_hight(self.tag_font_large_size, up_shape_width_if_array[1])  # width, hight
                                tag_shape_type = 'TAG_NORMAL'
                                tag_shape_left = shape_left + tag_up_offset_x
                                tag_shape_top = shape_top - shape_width_hight_array[1] * 0.5
                                tag_shape_width = shape_width_hight_array[0]
                                tag_shape_hight = shape_width_hight_array[1]
                                tag_shape_text = up_shape_width_if_array[1]

                                if action_type == 'CREATE':
                                    self.shape = self.slide.shapes
                                    ns_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,up_shape_width_if_array[0],shape_text])

                                # reflect description ip address name distance-1
                                flag_match_shape_width_if_array = False
                                for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                    if [shape_text, up_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                        tag_up_offset_x += tmp_shape_width_if_array[1]
                                        flag_match_shape_width_if_array = True
                                        tmp_add_width = tmp_shape_width_if_array[1]

                                if flag_match_shape_width_if_array == False:
                                    tag_up_offset_x += tag_shape_width + self.between_l3if
                                    tmp_add_width = tag_shape_width + self.between_l3if

                                #tag_up_offset_x += tag_shape_width + self.between_l3if

                                '''wp_right write ip address(up side)'''
                                offset_ipaddress = 0.0  # inches
                                remake_array = eval(tmp_update_l3_table_array[6])
                                if remake_array != []:
                                    for tmp_remake_array in remake_array:
                                        #print('##wp_right write ip address', tmp_remake_array)
                                        tag_shape_type = 'IP_ADDRESS_TAG'
                                        tag_ip_width = ns_def.get_description_width_hight(self.shae_font_size, tmp_remake_array[2])[0]
                                        #self.shape = self.slide.shapes
                                        #ns_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_tmp_update_l3_table_array[2])
                                        self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2], tmp_remake_array, shape_text, len(remake_array), [shape_text, up_shape_width_if_array[0]],tag_shape_left])

                                        offset_ipaddress += tag_shape_hight

                                        # reflect description ip address name distance-2
                                        #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                        #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width


                        # wite down side l3 if
                        for down_shape_width_if_array in self.shape_width_if_array[2]:
                            if down_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                shape_width_hight_array = ns_def.get_description_width_hight(self.tag_font_large_size, down_shape_width_if_array[1])  # width, hight
                                tag_shape_type = 'TAG_NORMAL'
                                tag_shape_left = shape_left + tag_down_offset_x
                                tag_shape_top = shape_top + shape_hight - shape_width_hight_array[1] * 0.5
                                tag_shape_width = shape_width_hight_array[0]
                                tag_shape_hight = shape_width_hight_array[1]
                                tag_shape_text = down_shape_width_if_array[1]

                                if action_type == 'CREATE':
                                    self.shape = self.slide.shapes
                                    ns_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,down_shape_width_if_array[0],shape_text])

                                # reflect description ip address name distance-1
                                flag_match_shape_width_if_array = False
                                for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                    if [shape_text, down_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                        tag_up_offset_x += tmp_shape_width_if_array[1]
                                        flag_match_shape_width_if_array = True
                                        tmp_add_width = tmp_shape_width_if_array[1]

                                if flag_match_shape_width_if_array == False:
                                    tag_down_offset_x += tag_shape_width + self.between_l3if
                                    tmp_add_width = tag_shape_width + self.between_l3if

                                #tag_down_offset_x += tag_shape_width + self.between_l3if

                                '''wp_right write ip address(down side)'''
                                offset_ipaddress = 0.0  # inches
                                remake_array = []
                                remake_array = eval(tmp_update_l3_table_array[6])
                                if remake_array != []:
                                    for tmp_remake_array in remake_array:
                                        #print('##wp_right write ip address', tmp_remake_array)
                                        tag_shape_type = 'IP_ADDRESS_TAG'
                                        tag_ip_width = ns_def.get_description_width_hight(self.shae_font_size, tmp_remake_array[2])[0]
                                        #self.shape = self.slide.shapes
                                        #ns_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top + tag_shape_hight + offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_tmp_update_l3_table_array[2])
                                        self.mark_multi_ip_array.append( [tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top + tag_shape_hight + offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2], tmp_remake_array, shape_text, len(remake_array), [shape_text, down_shape_width_if_array[0]],tag_shape_left])

                                        offset_ipaddress += tag_shape_hight

                                        # reflect description ip address name distance-2
                                        if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                            tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width

                offset_shape_right = offset_shape_right + shape_width + self.between_shape_column

        '''
        write folder line and outline 
        '''
        print('--- self.area_position_array ---  # shape_left, shape_top, shape_width, shape_hight, shape_text')
        #print(self.area_position_array)

        ### write folder
        folder_shape_type = 'FOLDER_NORMAL'
        folder_shape_left = self.area_position_array[0]
        folder_shape_top = self.area_position_array[1]
        folder_shape_width = self.area_position_array[2]
        folder_shape_hight = self.area_position_array[3]
        folder_shape_text = self.area_position_array[4]

        if action_type == 'CREATE':
            self.shape = self.slide.shapes
            ns_ddx_figure.extended.add_shape(self, folder_shape_type, folder_shape_left, folder_shape_top, folder_shape_width, folder_shape_hight, folder_shape_text)
            self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
            self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer

        ### write outline
        outline_shape_type = 'OUTLINE_NORMAL'
        outline_shape_left = self.outline_position_array[0]
        outline_shape_top = self.outline_position_array[1]
        outline_shape_width = self.outline_position_array[2]
        outline_shape_hight = self.outline_position_array[3]
        outline_shape_text = ''

        ### Kyuusai if len(target_position_shape_array) == 1
        outline_shape_hight += 0.2


        if action_type == 'CREATE':
            self.shape = self.slide.shapes
            ns_ddx_figure.extended.add_shape(self, outline_shape_type, outline_shape_left, outline_shape_top, outline_shape_width, outline_shape_hight, outline_shape_text)


        '''
        loop write l3segment
        '''
        all_l3_netowrk_list =[]
        all_l3segment_l3_netowrk_list = []
        self.all_written_if_line_array = []
        self.all_written_line_position_array = []
        multi_ip_address_num = 1
        for index_5,tmp_l3segment_line_array in enumerate(l3segment_line_array):
            #print(index_5,tmp_l3segment_line_array)
            if len(l3segment_line_array) > index_5 + 1:
                if l3segment_line_array[index_5][0][0] < l3segment_line_array[index_5 + 1][0][0]:
                    start_x = l3segment_line_array[index_5][0][0]
                else:
                    start_x = l3segment_line_array[index_5 + 1][0][0]

                if l3segment_line_array[index_5][0][1] > l3segment_line_array[index_5 + 1][0][1]:
                    end_x = l3segment_line_array[index_5][0][1]
                else:
                    end_x = l3segment_line_array[index_5 + 1][0][1]

                for index_55,tmp_tmp_l3segment_line_array in enumerate(tmp_l3segment_line_array[1]) :
                    '''pre check L3 SEGMENT'''
                    line_type = 'L3_SEGMENT'
                    inche_from_connect_x = start_x
                    inche_from_connect_y = tmp_tmp_l3segment_line_array
                    inche_to_connect_x = end_x
                    inche_to_connect_y = tmp_tmp_l3segment_line_array
                    #print(line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)
                    #ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                    shape_text = 'xxx.xxx.xxx.xxx/xx'
                    shape_type = 'L3_SEGMENT_GRAY'
                    shape_width = inche_to_connect_x - inche_from_connect_x
                    shape_hight = ns_def.get_description_width_hight(self.tag_font_large_size, shape_text)[1] * l3_segment_hight_ratio  # l3 segment hight ratio
                    shape_left = inche_from_connect_x
                    shape_top = inche_from_connect_y - shape_hight * 0.5

                    #self.shape = self.slide.shapes
                    #ns_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)
                    #self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
                    #self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer

                    '''write lines between device and l3 segment'''
                    # identify key l3 if for current broadcast domain
                    #print('##self.connected_l3if_key_array[index_5][index_55]  ', self.connected_l3if_key_array[index_5][index_55])

                    ###get l3IF of the broadcast domain

                    edge_left_x = 999.0   # inches
                    edge_right_x = -999.0  # inches

                    l3segment_edge_array = []
                    l3_network_list = []
                    the_l3segment_l3_network_list = []
                    for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
                        if self.connected_l3if_key_array[index_5][index_55] in tmp_target_l2_broadcast_group_array[1]:
                            #print('##MATCH  ',self.connected_l3if_key_array[index_5][index_55])
                            #print('##Target IFs  ', len(tmp_target_l2_broadcast_group_array[1]),tmp_target_l2_broadcast_group_array[1])
                            for tmp_all_l3if_tag_array in self.all_l3if_tag_array:
                                #print([tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]] , tmp_target_l2_broadcast_group_array[1])
                                if [tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]] in tmp_target_l2_broadcast_group_array[1]:
                                    #print('## Write line ', tmp_all_l3if_tag_array)

                                    line_type = 'L3_SEGMENT-L3IF'
                                    inche_from_connect_x = tmp_all_l3if_tag_array[1] + tmp_all_l3if_tag_array[3] * 0.5
                                    inche_from_connect_y = tmp_all_l3if_tag_array[2] + tmp_all_l3if_tag_array[4]
                                    inche_to_connect_x = tmp_all_l3if_tag_array[1] + tmp_all_l3if_tag_array[3] * 0.5
                                    inche_to_connect_y = tmp_tmp_l3segment_line_array - shape_hight * 0.5

                                    #for up
                                    if [tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]] in self.up_down_l3if_array[0]:
                                        inche_from_connect_y = tmp_all_l3if_tag_array[2]
                                        inche_to_connect_y = tmp_tmp_l3segment_line_array + shape_hight * 0.5

                                    '''check to exist the near vertical line. if there is near line, make between size to about 0.1 inches'''
                                    for tmp_all_written_line_position_array in self.all_written_line_position_array:
                                        if tmp_all_written_line_position_array[2] < inche_from_connect_y < tmp_all_written_line_position_array[4] or tmp_all_written_line_position_array[2] > inche_from_connect_y > tmp_all_written_line_position_array[4] \
                                                or tmp_all_written_line_position_array[2] < inche_to_connect_y < tmp_all_written_line_position_array[4] or tmp_all_written_line_position_array[2] > inche_to_connect_y > tmp_all_written_line_position_array[4]:
                                            if inche_from_connect_x - min_between_line <  tmp_all_written_line_position_array[1]   < inche_from_connect_x + min_between_line:
                                                if inche_from_connect_y != tmp_all_written_line_position_array[2] and inche_to_connect_y != tmp_all_written_line_position_array[2] and inche_from_connect_y != tmp_all_written_line_position_array[4] and inche_to_connect_y != tmp_all_written_line_position_array[4]:
                                                    if (inche_from_connect_x - tmp_all_written_line_position_array[1]) < 0:
                                                        offset_if_line = min_between_line + (inche_from_connect_x - tmp_all_written_line_position_array[1])
                                                    else:
                                                        offset_if_line = (inche_from_connect_x - tmp_all_written_line_position_array[1]) - min_between_line

                                                    inche_from_connect_x -= offset_if_line
                                                    inche_to_connect_x -= offset_if_line

                                    # print(line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)
                                    if action_type == 'CREATE':
                                        ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)
                                    self.all_written_line_position_array.append([line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y])

                                    if inche_from_connect_x < edge_left_x:
                                        edge_left_x = inche_from_connect_x
                                    if inche_from_connect_x > edge_right_x:
                                        edge_right_x = inche_from_connect_x

                                    #make l3 network address list
                                    for tmp_update_l3_table_array in self.update_l3_table_array:
                                        if tmp_update_l3_table_array[1] == tmp_all_l3if_tag_array[7] and tmp_update_l3_table_array[2] == tmp_all_l3if_tag_array[6]:
                                            remake_array = eval(tmp_update_l3_table_array[6])

                                            for tmp_remake_array in remake_array:
                                                l3_network_list.append(tmp_remake_array[1])
                                                all_l3_netowrk_list.append([tmp_remake_array[1],[tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]]])
                                                the_l3segment_l3_network_list.append([tmp_remake_array[1],[tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]]])


                                    #mark written if line for write ip address
                                    self.all_written_if_line_array.append([tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]])

                    l3segment_edge_array = [edge_left_x,edge_right_x ]
                    l3_network_list = ns_def.get_l2_broadcast_domains.get_unique_list(l3_network_list)

                    #print('--- l3_network_list ---')
                    #print(l3_network_list )



                    '''write l3 segment'''
                    #print('### l3segment_edge_array   ',l3segment_edge_array)

                    if len(l3_network_list) >= 2:
                        tmp_text = ''
                        for index_01, tmp_l3_network_list in enumerate(l3_network_list):
                            tmp_text += '(' + str(ns_def.num2alpha(multi_ip_address_num)) + ')' + tmp_l3_network_list + '  '

                            # pre match to l3if ip address
                            for index_11,tmp_the_l3segment_l3_network_list in enumerate(the_l3segment_l3_network_list):
                                if tmp_the_l3segment_l3_network_list[0] == tmp_l3_network_list:
                                    the_l3segment_l3_network_list[index_11].append('(' + str(ns_def.num2alpha(multi_ip_address_num)) + ')')

                            multi_ip_address_num += 1

                        shape_text = tmp_text


                    elif l3_network_list != []:
                        shape_text = str(l3_network_list[0])

                    else:
                        shape_text = ''


                    #print('--- the_l3segment_l3_network_list ---')
                    #print(the_l3segment_l3_network_list)

                    all_l3segment_l3_netowrk_list.extend(the_l3segment_l3_network_list)


                    shape_left = l3segment_edge_array[0]
                    shape_width = l3segment_edge_array[1] - l3segment_edge_array[0]

                    #check width text
                    tmp_text_width = ns_def.get_description_width_hight(self.tag_font_large_size, shape_text)[0]
                    if tmp_text_width > shape_width:
                        shape_width = tmp_text_width

                    #add left right inches
                    tmp_add_width = ns_def.get_description_width_hight(self.tag_font_large_size, 'aa')[0]
                    shape_left -= tmp_add_width
                    shape_width += tmp_add_width * 2

                    #check min width
                    tmp_char_width = ns_def.get_description_width_hight(self.tag_font_large_size,shape_text)[0]
                    if shape_width < tmp_char_width:
                        shape_width = tmp_char_width

                    #write l3 segment
                    if action_type == 'CREATE':
                        self.shape = self.slide.shapes
                        ns_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)
                        self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
                        self.slide.shapes._spTree.insert(3, self.shape._element)  # move shape to back layer


        #print('#### self.all_written_line_position_array  | line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y \n'  ,self.all_written_line_position_array , len(self.all_written_line_position_array))

        '''write ip on L3 IF'''
        #print('--- self.mark_multi_ip_array ---  [tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2],tmp_remake_array,shape_text ]')
        print('--- self.mark_multi_ip_array ---')
        #print(self.mark_multi_ip_array)
        print('--- self.all_written_if_line_array) ---')
        #print(self.all_written_if_line_array)
        print('--- all_l3segment_l3_netowrk_list ---')
        #print(all_l3segment_l3_netowrk_list)

        if self.mark_multi_ip_array != []:
            for tmp_mark_multi_ip_array in self.mark_multi_ip_array:
                if tmp_mark_multi_ip_array[8] == 1:
                    tag_ip_text = tmp_mark_multi_ip_array[5]
                    tag_ip_width = ns_def.get_description_width_hight(self.tag_font_large_size, tag_ip_text)[0]
                    tag_ip_text_2 = tmp_mark_multi_ip_array[6][0]
                    tag_ip_width_2 = ns_def.get_description_width_hight(self.tag_font_large_size, tag_ip_text_2 + str('x'))[0]
                    tag_ip_left_2 = tmp_mark_multi_ip_array[10]

                    #chack add (x)
                    for tmp_all_l3segment_l3_netowrk_list in all_l3segment_l3_netowrk_list:
                        if tmp_all_l3segment_l3_netowrk_list[1] == tmp_mark_multi_ip_array[9] and tmp_all_l3segment_l3_netowrk_list[0] == tmp_mark_multi_ip_array[6][1] and len(tmp_all_l3segment_l3_netowrk_list) == 3:
                            tag_ip_text = tmp_all_l3segment_l3_netowrk_list[2] + tmp_mark_multi_ip_array[5]
                            tag_ip_width = ns_def.get_description_width_hight(self.tag_font_large_size, tag_ip_text)[0]
                            tag_ip_text_2 = tmp_all_l3segment_l3_netowrk_list[2] + tmp_mark_multi_ip_array[6][0]
                            tag_ip_width_2 = ns_def.get_description_width_hight(self.tag_font_large_size, tag_ip_text_2 + str('x'))[0]

                    # check no line l3 if
                    if tmp_mark_multi_ip_array[9] in self.all_written_if_line_array:
                        if action_type == 'CREATE':
                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, tmp_mark_multi_ip_array[0], tmp_mark_multi_ip_array[1], tmp_mark_multi_ip_array[2], tag_ip_width, tmp_mark_multi_ip_array[4], tag_ip_text)
                    else:
                        if action_type == 'CREATE':
                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, tmp_mark_multi_ip_array[0], tag_ip_left_2, tmp_mark_multi_ip_array[2], tag_ip_width_2, tmp_mark_multi_ip_array[4], tag_ip_text_2)

                else:
                    for tmp_all_l3segment_l3_netowrk_list in all_l3segment_l3_netowrk_list:
                        if tmp_all_l3segment_l3_netowrk_list[1] == tmp_mark_multi_ip_array[9] and tmp_all_l3segment_l3_netowrk_list[0] == tmp_mark_multi_ip_array[6][1]:
                            tag_ip_text = tmp_all_l3segment_l3_netowrk_list[2] + tmp_mark_multi_ip_array[5]
                            tag_ip_width = ns_def.get_description_width_hight(self.tag_font_large_size, tag_ip_text)[0]
                            tag_ip_text_2 = tmp_all_l3segment_l3_netowrk_list[2]  + tmp_mark_multi_ip_array[6][0]
                            tag_ip_width_2 = ns_def.get_description_width_hight(self.tag_font_large_size, tag_ip_text_2 + str('x'))[0]
                            tag_ip_left_2 = tmp_mark_multi_ip_array[10]

                            # print('### tag_ip_text, tag_ip_width',tag_ip_text,tag_ip_width)

                            # check no line l3 if
                            if tmp_mark_multi_ip_array[9] in self.all_written_if_line_array:
                                if action_type == 'CREATE':
                                    self.shape = self.slide.shapes
                                    ns_ddx_figure.extended.add_shape(self, tmp_mark_multi_ip_array[0], tmp_mark_multi_ip_array[1], tmp_mark_multi_ip_array[2], tag_ip_width, tmp_mark_multi_ip_array[4], tag_ip_text)
                            else:
                                if action_type == 'CREATE':
                                    self.shape = self.slide.shapes
                                    ns_ddx_figure.extended.add_shape(self, tmp_mark_multi_ip_array[0], tag_ip_left_2, tmp_mark_multi_ip_array[2], tag_ip_width_2, tmp_mark_multi_ip_array[4], tag_ip_text_2)
                            break
        '''
        Write line of L3 instance
        '''
        if self.size_l3_instance_array != []:
            print('--- self.size_l3_instance_array ---  shape_text, l3_shape_text, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight')
            #print(self.size_l3_instance_array)
            print('--- self.all_l3if_tag_array --- , tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,up_shape_width_if_array[0],shape_text')
            #print(self.all_l3if_tag_array)

        used_line_array = []
        for tmp_update_l3_instance_array in self.update_l3_instance_array:
            for tmp_all_l3if_tag_array in self.all_l3if_tag_array:
                if tmp_update_l3_instance_array[0][1] == tmp_all_l3if_tag_array[7] and tmp_update_l3_instance_array[0][2] == tmp_all_l3if_tag_array[6]:
                    line_type = 'L3_INSTANCE'
                    inche_from_connect_x = tmp_all_l3if_tag_array[1] + tmp_all_l3if_tag_array[3] * 0.5
                    inche_from_connect_y = tmp_all_l3if_tag_array[2]

                    for tmp_size_l3_instance_array in self.size_l3_instance_array:
                        if tmp_size_l3_instance_array[0] == tmp_update_l3_instance_array[0][1] and tmp_size_l3_instance_array[1] == tmp_update_l3_instance_array[1]:
                            inche_to_connect_x = tmp_size_l3_instance_array[3] + tmp_size_l3_instance_array[5] * 0.5
                            inche_to_connect_y = tmp_size_l3_instance_array[4] + tmp_size_l3_instance_array[6]

                            # for up
                            if [tmp_all_l3if_tag_array[7], tmp_all_l3if_tag_array[6]] in self.up_down_l3if_array[0]:
                                inche_from_connect_y = tmp_all_l3if_tag_array[2] + tmp_all_l3if_tag_array[4]
                                inche_to_connect_y = tmp_size_l3_instance_array[4]

                            if action_type == 'CREATE' and [line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y] not in used_line_array:
                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)
                                used_line_array.append([line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y])
                                break

        return ([outline_shape_type, outline_shape_left, outline_shape_top, outline_shape_width, outline_shape_hight, folder_shape_text])




'''
LOCAL DEF
'''

def get_shape_width_if_array(self,device_name):
    #get if array
    tmp_up_array = []
    tmp_down_array = []
    full_ip_address_width_array = []
    for tmp_up_down_l3if_array in self.up_down_l3if_array[0]:
        if device_name == tmp_up_down_l3if_array[0]:
            tmp_return = ns_def.get_tag_name_from_full_name(device_name, tmp_up_down_l3if_array[1], self.position_line_tuple)
            if tmp_return == '_NO_MATCH_':
                tmp_if_array = ns_def.adjust_portname(tmp_up_down_l3if_array[1])
                tmp_up_array.append([tmp_up_down_l3if_array[1],str(tmp_if_array[0] + ' ' + tmp_if_array[2]).replace('  ',' ')])
            else:
                tmp_up_array.append([tmp_up_down_l3if_array[1],tmp_return])

    for tmp_up_down_l3if_array in self.up_down_l3if_array[1]:
        if device_name == tmp_up_down_l3if_array[0]:
            tmp_return = ns_def.get_tag_name_from_full_name(device_name, tmp_up_down_l3if_array[1], self.position_line_tuple)
            if tmp_return == '_NO_MATCH_':
                tmp_if_array = ns_def.adjust_portname(tmp_up_down_l3if_array[1])
                tmp_down_array.append([tmp_up_down_l3if_array[1],str(tmp_if_array[0] + ' ' + tmp_if_array[2]).replace('  ',' ')])
            else:
                tmp_down_array.append([tmp_up_down_l3if_array[1],tmp_return])

    #calc width for device
    tmp_up_width = self.between_l3if  #inches
    tmp_down_width = self.between_l3if   # inches
    full_ip_address = 'xxx.xxx.xxx.xxx/xxxx'
    distance_full_ip_address = ns_def.get_description_width_hight(self.tag_font_large_size,full_ip_address)[0]


    for tmp_tmp_up_array in tmp_up_array:
        shape_width_hight_array = ns_def.get_description_width_hight(self.tag_font_large_size,tmp_tmp_up_array[1])
        first_up_width = shape_width_hight_array[0] + self.between_l3if

        #check distance of full ip address and not connect l3segment
        tmp_tmp_up_width = distance_full_ip_address
        if tmp_tmp_up_width > first_up_width and [device_name,tmp_tmp_up_array[0]] not in self.l3_if_has_l3_segment_array:
            first_up_width = tmp_tmp_up_width
            full_ip_address_width_array.append([[device_name,tmp_tmp_up_array[0]],tmp_tmp_up_width])

        tmp_up_width += first_up_width


    for tmp_tmp_down_array in tmp_down_array:
        shape_width_hight_array = ns_def.get_description_width_hight(self.tag_font_large_size,tmp_tmp_down_array[1])
        first_down_width =  shape_width_hight_array[0] + self.between_l3if

        # check distance of full ip address
        tmp_tmp_down_width = distance_full_ip_address
        if tmp_tmp_down_width > first_down_width and [device_name,tmp_tmp_down_array[0]]  not in self.l3_if_has_l3_segment_array:
            first_down_width = tmp_tmp_down_width
            full_ip_address_width_array.append([[device_name, tmp_tmp_down_array[0]], tmp_tmp_down_width])

        tmp_down_width += first_down_width


    if tmp_up_width > tmp_down_width:
        return_shape_width = tmp_up_width
    else:
        return_shape_width = tmp_down_width

    #print('### return_shape_width, tmp_up_array,tmp_down_array',return_shape_width,tmp_up_array,tmp_down_array,full_ip_address_width_array])
    return([return_shape_width,tmp_up_array,tmp_down_array,full_ip_address_width_array])

def get_up_down_l3if_count(self,marge_target_position_shape_array):
    buttom_array = []
    top_array = []
    for index_8,tmp_marge_target_position_shape_array in enumerate(marge_target_position_shape_array):
        for tmp_tmp_marge_target_position_shape_array in tmp_marge_target_position_shape_array:

            if '_AIR_' not in tmp_tmp_marge_target_position_shape_array:
                #print(tmp_tmp_marge_target_position_shape_array)
                tmp_buttom_array = []
                tmp_top_array = []

                for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
                    for tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:

                        if tmp_tmp_target_l2_broadcast_group_array[0] == tmp_tmp_marge_target_position_shape_array:
                            #print(tmp_tmp_target_l2_broadcast_group_array)

                            #check if l3if connect to buttom side l3segment
                            flag_buttom_exist = False
                            if len(marge_target_position_shape_array)  > index_8:
                                for i in range(index_8 + 1,len(marge_target_position_shape_array)):

                                    for buttom_marge_target_position_shape_array in marge_target_position_shape_array[i]:
                                        #print(buttom_marge_target_position_shape_array)
                                        for now_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                                            if buttom_marge_target_position_shape_array == now_tmp_target_l2_broadcast_group_array[0]:
                                                #print('## Buttom l3 if ',now_tmp_target_l2_broadcast_group_array)
                                                flag_buttom_exist = True

                            #check if l3if connect to top side l3segment
                            flag_top_exist = False
                            if 1  <= index_8:
                                for i in range(0,index_8):

                                    for top_marge_target_position_shape_array in marge_target_position_shape_array[i]:
                                        #print(buttom_marge_target_position_shape_array)
                                        for now_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                                            if top_marge_target_position_shape_array == now_tmp_target_l2_broadcast_group_array[0]:
                                                #print('## Top l3 if ',now_tmp_target_l2_broadcast_group_array)
                                                flag_top_exist = True

                            #print('### flag_buttom_exist,flag_top_exist  ',flag_buttom_exist,flag_top_exist,tmp_tmp_target_l2_broadcast_group_array)
                            if self.new_wp_exist_array[0] == [] and index_8 == 0:
                                tmp_buttom_array.append(tmp_tmp_target_l2_broadcast_group_array)

                            elif flag_buttom_exist == False and flag_top_exist == False:
                                tmp_top_array.append(tmp_tmp_target_l2_broadcast_group_array)

                            elif flag_buttom_exist == True and flag_top_exist == False:
                                tmp_buttom_array.append(tmp_tmp_target_l2_broadcast_group_array)

                            elif flag_buttom_exist == False and flag_top_exist == True:
                                tmp_top_array.append(tmp_tmp_target_l2_broadcast_group_array)

                            elif flag_buttom_exist == True and flag_top_exist == True:
                                tmp_top_array.append(tmp_tmp_target_l2_broadcast_group_array)

                #print('### UP /DOWN  ',tmp_tmp_marge_target_position_shape_array,tmp_top_array,tmp_buttom_array)
                buttom_array.extend(tmp_buttom_array)
                top_array.extend(tmp_top_array)

    buttom_array = ns_def.get_l2_broadcast_domains.get_unique_list(buttom_array)
    top_array = ns_def.get_l2_broadcast_domains.get_unique_list(top_array)
    return ([top_array,buttom_array])


def check_move_to_right(self,top_device_name_array,target_position_shape_array):
    #print('--- check_move_to_right ---')
    if  self.index_2 >= 1 and self.index_2 <= len(self.index_1_array) - 2:
        for i in range(self.index_2):
            #print(target_position_shape_array[self.index_2] ,target_position_shape_array[self.index_2 + 1],target_position_shape_array[i], '  #original ,source , target,') # source , target
            source_array = target_position_shape_array[self.index_2+ 1]
            target_array = target_position_shape_array[i]

            '''for tmp_new_wp_exist_array_3 in self.new_wp_exist_array[3]:
                if tmp_new_wp_exist_array_3 in source_array:
                    source_array.remove(tmp_new_wp_exist_array_3)
                if tmp_new_wp_exist_array_3 in target_array:
                    target_array.remove(tmp_new_wp_exist_array_3)'''

            for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
                flag_source = False
                flag_target = False
                for tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                    if tmp_tmp_target_l2_broadcast_group_array != []:
                        #print(tmp_tmp_target_l2_broadcast_group_array,source_array,target_array)
                        if tmp_tmp_target_l2_broadcast_group_array[0] in source_array:
                            flag_source = True
                        if tmp_tmp_target_l2_broadcast_group_array[0] in target_array:
                            flag_target = True


                if flag_source == True and flag_target == True:
                    #print(flag_source ,flag_target )
                    return (True)

    return (False)

def get_l3_segment_num(self,top_device_name_array,target_position_shape_array):
    #print('--- get_l3_segment_num ---')
    #print(top_device_name_array )
    count_l3segment = 0
    connected_l3if_key_array = []
    tmp_used_l3segment_array = []

    #segment that connected to both top and buttom
    if  self.index_2 + 1 in self.index_1_array:
        #print('--- self.index_2,self.index_1_array ---',self.index_2,self.index_1_array)
        for i in range(self.index_2 + 1 , len(self.index_1_array)):
            buttom_device_name_array = target_position_shape_array[i]
            for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
                if tmp_target_l2_broadcast_group_array[0] not in self.used_l3segment_array:
                    for tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                        if tmp_tmp_target_l2_broadcast_group_array[0] in top_device_name_array and tmp_tmp_target_l2_broadcast_group_array not in connected_l3if_key_array:
                            #print('tmp_tmp_target_l2_broadcast_group_array[0],top_device_name_array', tmp_tmp_target_l2_broadcast_group_array[0], top_device_name_array)
                            for tmp_tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                                if tmp_tmp_tmp_target_l2_broadcast_group_array[0] in buttom_device_name_array:
                                    count_l3segment += 1
                                    connected_l3if_key_array.append(tmp_tmp_target_l2_broadcast_group_array )
                                    tmp_used_l3segment_array.append(tmp_target_l2_broadcast_group_array[0])
                                    break

        self.used_l3segment_array.extend(tmp_used_l3segment_array)

    #segment that connected to left or right only. target is buttom side shape
    #segment that is not connected to top and buttom. target is buttom side shape

    if  self.index_2 + 1 in self.index_1_array:
        buttom_device_name_array = target_position_shape_array[self.index_2 + 1]
        #print('### buttom_device_name_array ',buttom_device_name_array)
        for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
            tmp_count = 0
            for tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                if tmp_tmp_target_l2_broadcast_group_array[0] in buttom_device_name_array:
                    tmp_count += 1

            if tmp_count == len(tmp_target_l2_broadcast_group_array[1]) and tmp_count >= 2:
                #print('### tmp_count, tmp_target_l2_broadcast_group_array[1] ',tmp_count, tmp_target_l2_broadcast_group_array[1])
                count_l3segment += 1
                connected_l3if_key_array.append(tmp_tmp_target_l2_broadcast_group_array)

            #if tmp_count == len(tmp_target_l2_broadcast_group_array[1]) and tmp_count == 1:
                #print('*keep as loopback etc..  ' , tmp_count, tmp_target_l2_broadcast_group_array[1])


    ### Kyuusai , there is not upside wp and firest device row (self.index_2 = 0 )
    if self.index_2 == 0 and self.new_wp_exist_array[0] == []:
        for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
            tmp_count = 0
            if tmp_target_l2_broadcast_group_array[0] not in self.used_l3segment_array:
                for tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                    if tmp_tmp_target_l2_broadcast_group_array[0] in top_device_name_array and tmp_tmp_target_l2_broadcast_group_array not in connected_l3if_key_array:
                        # print('tmp_tmp_target_l2_broadcast_group_array[0],top_device_name_array', tmp_tmp_target_l2_broadcast_group_array[0], top_device_name_array)
                        for tmp_tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                            if tmp_tmp_tmp_target_l2_broadcast_group_array[0] in target_position_shape_array[0] and tmp_target_l2_broadcast_group_array[0] not in tmp_used_l3segment_array:
                                tmp_count += 1

                                if tmp_count >= 2:
                                    count_l3segment += 1
                                    connected_l3if_key_array.insert(0,tmp_tmp_target_l2_broadcast_group_array)
                                    tmp_used_l3segment_array.append(tmp_target_l2_broadcast_group_array[0])

        self.used_l3segment_array.extend(tmp_used_l3segment_array)

    #print('-- connected_l3if_key_array --', connected_l3if_key_array)
    return ([count_l3segment,connected_l3if_key_array])





